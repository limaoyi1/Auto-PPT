import warnings

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Pt, Inches

from cv import get_largest_white_region_in_slide
from picture import search, check_same_name_file, get_random_file, get_image_resolution

warnings.filterwarnings('ignore')


def make_page_left(presentation, theme, numbers, titles, texts, keywords, index):
    """
    在左边生成文本框，右边生成图片
    :param presentation:
    :param numbers:
    :param titles:
    :param texts:
    :param keywords:
    :param index:
    :return:
    """
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    # 1、新增母版占位符
    placeholder1 = slide.placeholders[1]
    path = get_random_file(theme)
    picture = placeholder1.insert_picture(path)
    placeholder2 = slide.placeholders[2]
    placeholder2.element.getparent().remove(placeholder2.element)
    # 2、设置占位符宽高
    picture.left = 0
    picture.top = 0
    picture.width = presentation.slide_width
    picture.height = presentation.slide_height
    # 3、添加文本框1
    shapes = slide.shapes
    text_box = shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(3), Inches(0.8))
    tf = text_box.text_frame  # TextFrame
    tf.clear()  # Clear existing content
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # 添加标题
    paragraph = tf.paragraphs[0]
    text = str(numbers[index]) + " " + titles[index]
    paragraph.text = text
    paragraph.font.bold = True
    paragraph.font.name = "微软雅黑"
    paragraph.font.color.rgb = RGBColor(51, 0, 102)
    paragraph.font.size = Pt(26)
    paragraph.word_wrap = True
    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # 3、添加 shape 集合
    subtitle_by_index = get_subtitle_by_index(numbers, index)
    print(subtitle_by_index, '=================================')
    group_shape = slide.shapes.add_group_shape()
    group_shape.height = presentation.slide_height
    # 偏移值
    offset = 0
    for idx in range(len(subtitle_by_index)):
        if idx != 0:
            text_box_common = group_shape.shapes.add_textbox(Inches(0.8), Inches(1.1 + offset), Inches(4.8),
                                                             Inches(0.4))
            offset = offset + 0.4
            text_frame_common = text_box_common.text_frame
            text_frame_common.paragraphs[0].text = titles[subtitle_by_index[idx]]
            text_frame_common.paragraphs[0].font.bold = True
            text_frame_common.paragraphs[0].font.name = "微软雅黑"
            text_frame_common.paragraphs[0].font.color.rgb = RGBColor(0, 0, 102)
            text_frame_common.paragraphs[0].font.size = Inches(0.25)
            text_frame_common.paragraphs[0].font.word_wrap = True
            text_frame_common.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            if texts[subtitle_by_index[idx]] is not None and str(texts[subtitle_by_index[idx]]) != "nan":
                height = (len("\t" + texts[subtitle_by_index[idx]]) / (4.8 / 0.15) + 1) * 0.16 + 0.1
                text_box_common2 = group_shape.shapes.add_textbox(Inches(0.7), Inches(1.1 + offset), Inches(4.8),
                                                                  Inches(height))
                offset = height + offset + 0.15
                text_frame_common2 = text_box_common2.text_frame
                text_frame_common2.paragraphs[0].text = str("\t" + texts[subtitle_by_index[idx]])
                text_frame_common2.paragraphs[0].font.name = "微软雅黑"
                text_frame_common2.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                text_frame_common2.paragraphs[0].font.size = Inches(0.15)
                text_frame_common2.paragraphs[0].font.word_wrap = True
                text_frame_common2.word_wrap = True
                text_frame_common2.paragraphs[0].space_after = Inches(0.1)
                text_frame_common2.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                text_frame_common2.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # 添加插图
    # 获取二级插图keywords
    keywords_index_ = keywords[index]
    if keywords_index_ is None and keywords_index_ == "":
        pass
    else:
        b = check_same_name_file("./picture", str(keywords_index_) + ".jpg")
        if (b == True):
            link = "./picture" + "/" + str(keywords_index_) + ".jpg"
        else:
            link = search(str(keywords_index_))
        width1, height1 = get_image_resolution(link)
        left = Inches(5.5)  # 图片左侧距离幻灯片左边缘的距离
        top = Inches(1.2)  # 图片顶部距离幻灯片上边缘的距离
        width = Inches(4)  # 图片的宽度

        height = width * (height1 / width1)  # 图片的高度
        ph_picture = slide.shapes.add_picture(link, left, top, width, height)


def make_page_right(presentation, theme, numbers, titles, texts, keywords, index):
    """
    在左边生成文本框，右边生成图片
    :param presentation:
    :param numbers:
    :param titles:
    :param texts:
    :param keywords:
    :param index:
    :return:
    """
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    # 1、新增母版占位符
    placeholder1 = slide.placeholders[1]
    path = get_random_file(theme)
    picture = placeholder1.insert_picture(path)
    placeholder2 = slide.placeholders[2]
    placeholder2.element.getparent().remove(placeholder2.element)
    # 2、设置占位符宽高
    picture.left = 0
    picture.top = 0
    picture.width = presentation.slide_width
    picture.height = presentation.slide_height
    # 3、添加文本框1
    shapes = slide.shapes
    text_box = shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(3), Inches(0.8))
    tf = text_box.text_frame  # TextFrame
    tf.clear()  # Clear existing content
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # 添加标题
    paragraph = tf.paragraphs[0]
    text = str(numbers[index]) + " " + titles[index]
    paragraph.text = text
    paragraph.font.bold = True
    paragraph.font.name = "微软雅黑"
    paragraph.font.color.rgb = RGBColor(51, 0, 102)
    paragraph.font.size = Pt(26)
    paragraph.word_wrap = True
    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # 3、添加 shape 集合
    subtitle_by_index = get_subtitle_by_index(numbers, index)
    print(subtitle_by_index, '=================================')
    group_shape = slide.shapes.add_group_shape()
    group_shape.height = presentation.slide_height

    # 偏移值
    offset = 0
    for idx in range(len(subtitle_by_index)):
        if (idx != 0):
            text_box_common = group_shape.shapes.add_textbox(Inches(4.6), Inches(1.1 + offset), Inches(4.8),
                                                             Inches(0.4))
            offset = offset + 0.4
            # 设置文本框自适应
            text_box_common.word_wrap = True
            text_box_common.auto_size = True
            text_frame_common = text_box_common.text_frame
            text_frame_common.paragraphs[0].text = titles[subtitle_by_index[idx]]
            text_frame_common.paragraphs[0].font.bold = True
            text_frame_common.paragraphs[0].font.name = "微软雅黑"
            text_frame_common.paragraphs[0].font.color.rgb = RGBColor(0, 0, 102)
            text_frame_common.paragraphs[0].font.size = Inches(0.25)
            text_frame_common.paragraphs[0].font.word_wrap = True
            text_frame_common.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            text_frame_common.auto_size = MSO_AUTO_SIZE.NONE
            # cy = None
            # for child in text_box_common._element.spPr:
            #     print(child.tag)
            #     xfrm_attrs = child.attrib
            #     for sub_child in child:
            #         if 'off' in sub_child.tag:
            #             pass
            #             # off_attrs = sub_child.attrib
            #             # for attr_name, attr_value in off_attrs.items():
            #                 # print(f"off - {attr_name}: {attr_value}")
            #         elif 'ext' in sub_child.tag:
            #             ext_attrs = sub_child.attrib
            #             for attr_name, attr_value in ext_attrs.items():
            #                 # print(f"ext - {attr_name}: {attr_value}")
            #                 if attr_name == "cy":
            #                     cy = int(attr_value)
            # print(f"cy 等于 {cy}")
            if texts[subtitle_by_index[idx]] is not None and str(texts[subtitle_by_index[idx]]) != "nan":
                height = (len("\t" + texts[subtitle_by_index[idx]]) / (4.8 / 0.15) + 1) * 0.16 + 0.1
                text_box_common2 = group_shape.shapes.add_textbox(Inches(4.6), Inches(1.1 + offset), Inches(4.8),
                                                                  Inches(height))
                offset = height + offset + 0.15
                # 设置文本框自适应
                text_frame_common2 = text_box_common2.text_frame
                text_frame_common2.paragraphs[0].text = str("\t" + texts[subtitle_by_index[idx]])
                text_frame_common2.paragraphs[0].font.name = "微软雅黑"
                text_frame_common2.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                text_frame_common2.paragraphs[0].font.size = Inches(0.15)
                text_frame_common2.paragraphs[0].font.word_wrap = True
                text_frame_common2.word_wrap = True
                text_frame_common2.paragraphs[0].space_after = Inches(0.1)
                text_frame_common2.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                text_frame_common2.auto_size = MSO_AUTO_SIZE.NONE
    # 添加插图
    # 获取二级插图keywords
    keywords_index_ = keywords[index]
    if keywords_index_ is None and keywords_index_ == "":
        pass
    else:
        b = check_same_name_file("./picture", str(keywords_index_) + ".jpg")
        if (b == True):
            link = "./picture" + "/" + str(keywords_index_) + ".jpg"
        else:
            link = search(str(keywords_index_))
        width1, height1 = get_image_resolution(link)
        left = Inches(0.5)  # 图片左侧距离幻灯片左边缘的距离
        top = Inches(1.2)  # 图片顶部距离幻灯片上边缘的距离
        width = Inches(4)  # 图片的宽度

        height = width * (height1 / width1)  # 图片的高度
        ph_picture = slide.shapes.add_picture(link, left, top, width, height)


def get_subtitle_by_index(numbers, index):
    numbers_index_ = numbers[index]
    subtitle_index = []
    for idx in range(len(numbers)):
        if str(numbers[idx]).startswith(str(numbers_index_)):
            subtitle_index.append(idx)
    return subtitle_index


def make_page_cv(presentation, theme, numbers, titles, texts, keywords, index):
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    # 1、新增母版占位符
    placeholder1 = slide.placeholders[1]
    path = get_random_file(theme)
    # 通过 cv 获取 box 信息
    box_x, box_y, box_w, box_h = get_largest_white_region_in_slide(path)
    if get_largest_white_region_in_slide(path) == None or box_w < 1:
        make_page_right(presentation, theme, numbers, titles, texts, keywords, index)
        return
    picture = placeholder1.insert_picture(path)
    placeholder2 = slide.placeholders[2]
    placeholder2.element.getparent().remove(placeholder2.element)
    # 2、设置占位符宽高
    picture.left = 0
    picture.top = 0
    picture.width = presentation.slide_width
    picture.height = presentation.slide_height
    # 3、添加文本框1
    shapes = slide.shapes
    text_box = shapes.add_textbox(Inches(box_x), Inches(box_y), Inches(box_w), Inches(0.8))
    tf = text_box.text_frame  # TextFrame
    tf.clear()  # Clear existing content
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # 添加标题
    paragraph = tf.paragraphs[0]
    text = str(numbers[index]) + " " + titles[index]
    paragraph.text = text
    paragraph.font.bold = True
    paragraph.font.name = "微软雅黑"
    paragraph.font.color.rgb = RGBColor(51, 0, 102)
    paragraph.font.size = Pt(26)
    paragraph.word_wrap = True
    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # 3、添加 shape 集合
    subtitle_by_index = get_subtitle_by_index(numbers, index)
    print(subtitle_by_index, '=================================')
    group_shape = slide.shapes.add_group_shape()
    group_shape.height = presentation.slide_height

    for idx in range(len(subtitle_by_index)):
        if (idx != 0):
            text_box_common = group_shape.shapes.add_textbox(Inches(box_x) + Inches(0.25),
                                                             Inches(box_y) + Inches(0.4 + idx * 0.8),
                                                             Inches(box_w) - Inches(0.5),
                                                             Inches(0.8))
            text_frame_common = text_box_common.text_frame
            text_frame_common.paragraphs[0].text = titles[subtitle_by_index[idx]]
            text_frame_common.paragraphs[0].font.bold = True
            text_frame_common.paragraphs[0].font.name = "微软雅黑"
            text_frame_common.paragraphs[0].font.color.rgb = RGBColor(0, 0, 102)
            text_frame_common.paragraphs[0].font.size = Inches(0.25)
            text_frame_common.paragraphs[0].font.word_wrap = True
            text_frame_common.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            if texts[subtitle_by_index[idx]] is not None and str(texts[subtitle_by_index[idx]]) != "nan":
                text_box_common2 = group_shape.shapes.add_textbox(Inches(box_x) + Inches(0.25),
                                                                  Inches(box_y) + Inches(0.3 + idx * 0.8 + 0.5),
                                                                  Inches(box_w) - Inches(0.5),
                                                                  Inches(0.8))
                text_frame_common2 = text_box_common2.text_frame
                text_frame_common2.paragraphs[0].text = str("\t" + texts[subtitle_by_index[idx]])
                text_frame_common2.paragraphs[0].font.name = "微软雅黑"
                text_frame_common2.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                text_frame_common2.paragraphs[0].font.size = Inches(0.15)
                text_frame_common2.paragraphs[0].font.word_wrap = True
                text_frame_common2.word_wrap = True
                text_frame_common2.paragraphs[0].space_after = Inches(0.1)
                text_frame_common2.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                text_frame_common2.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
