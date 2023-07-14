import datetime
import os
from enum import Enum

from PIL.ImageQt import rgb
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.text.text import Font

from mdtree.parser import parse_string, Out, Heading
from mdtree.utils import get_random_theme, get_random_file, read_md_file


class Tree2PPT:
    prs: Presentation = None
    md_str: str = None
    out: Out = None
    tree: Heading = None
    theme: str = None

    def __init__(self, md_str1):
        self.init_pptx()
        self.init_markdown(md_str1)
        self.traverse_tree(self.tree)
        now = datetime.datetime.now().timestamp()
        path = './../myppt/test' + str(now) + '.pptx'
        if not os.path.exists('./../myppt'):
            os.makedirs('./../myppt')
        self.prs.save(path)
        pass

    def init_pptx(self):
        prs = Presentation()
        self.theme = get_random_theme()
        self.prs = prs

    def init_markdown(self, md_str):
        self.md_str = md_str
        self.out = parse_string(md_str)
        self.tree = self.out.main

    def traverse_tree(self, heading):
        if heading.source is None or heading.source == '':
            content = ""
            if heading.children is not []:
                for child in heading.children:
                    content = content + child.text + "\n"
            MD2Slide(self.prs, self.theme, heading.text, content=content)
        else:
            MD2Slide(self.prs, self.theme, heading.text, content=heading.source)

        # self.make_slide_demo(self.prs, heading.text, heading.source)
        MD2Slide(self.prs, self.theme, heading.text, content=heading.source)
        if heading.children is not []:
            for child in heading.children:
                self.traverse_tree(child)


class MarkdownCategory:
    TITLE = "#"
    CONTENT = "<p>"

    pass


class MD2Slide:
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "微软雅黑"
    font_title_size: Pt = Pt(26)
    font_content_size: Pt = Pt(18)
    font_title_color: rgb = RGBColor(51, 0, 102)
    font_content_color: rgb = RGBColor(51, 0, 102)

    def __init__(self, presentation, theme_path, title, content, *args, **kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[8])
        self.title = title
        self.content = content
        self.theme = theme_path
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()

    def init_slide(self):
        placeholder1 = self.slide.placeholders[1]
        path = get_random_file(self.theme)
        picture = placeholder1.insert_picture(path)
        placeholder2 = self.slide.placeholders[2]
        placeholder2.element.getparent().remove(placeholder2.element)
        # 2、设置占位符宽高
        picture.left = 0
        picture.top = 0
        picture.width = self.presentation.slide_width
        picture.height = self.presentation.slide_height

    def init_font(self, **kwargs):
        if 'font_name' in kwargs:
            self.font_name = kwargs['font_name']
        if 'font_title_size' in kwargs:
            self.font_title_size = kwargs['font_title_size']
        if 'font_content_size' in kwargs:
            self.font_content_size = kwargs['font_content_size']
        if 'font_title_color' in kwargs:
            self.font_title_color = kwargs['font_title_color']
        if 'font_content_color' in kwargs:
            self.font_content_color = kwargs['font_content_color']

    def get_font(self, font: Font, category: str):
        font.bold = True
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(3), Inches(0.8))
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # 添加标题
        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def init_content(self):
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        tf = text_box_content.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True
        # 添加正文
        paragraph = tf.paragraphs[0]
        paragraph.text = self.content
        self.get_font(paragraph.font, MarkdownCategory.CONTENT)
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


if __name__ == "__main__":
    md_content = read_md_file("./txt.md")
    out = parse_string(md_content)
    Tree2PPT(md_content)
