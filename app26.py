import os

from langchain.llms import OpenAI
from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT, PP_ALIGN
import markdown_to_json

import codecs
import configparser
config = configparser.ConfigParser()
# 指定编码为 UTF-8
config.read_file(codecs.open('config1.ini', 'r', 'utf-8-sig'))
OPENAI_API_KEY = config.get('Credentials', 'OPENAI_API_KEY')

os.environ["OPENAI_API_KEY"] = OPENAI_API_KEY
prs = Presentation()


def get_gpt_data():
    llm = OpenAI(temperature=0.9, model_name="gpt-3.5-turbo", max_tokens=4000)
    text = "写一份PPT，要求有层次的大纲，含有二级目录，目录不带序号，主题是新人如何做好直播，要求每个大纲下面有详细内容，不需要总结，不少于2000字,用markdown呈现,呈现格式必须是markdown"
    return llm(text)


def add_title(index, item):
    print(index, item)


# markdown格式是多级树结构
# 第一级是 # 区分，第二级是 ## 区分， 第三极是 ### 区分,...
# 此处用markdown_to_json 转化为标准json格式，方便后续处理
def clean_data(chat_str):
    print(chat_str)
    return markdown_to_json.dictify(chat_str)


def render_item(key, value):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    # 1、新增母版占位符
    placeholder1 = slide.placeholders[1]
    placeholder2 = slide.placeholders[2]
    picture = placeholder1.insert_picture('./pptx/static/bg.png')
    placeholder2.element.getparent().remove(placeholder2.element)

    # 2、设置占位符宽高
    picture.left = 0
    picture.top = 0
    picture.width = prs.slide_width
    picture.height = prs.slide_height
    # 3、添加文本框
    text_box = slide.shapes.add_textbox(Cm(5), Cm(0), Cm(15), Cm(19))
    tf = text_box.text_frame
    tf.text = key
    tf.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
    tf.paragraphs[0].font.size = Pt(46)
    tf.paragraphs[0].font.bold = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    # tf.paragraphs[0].alignment = MSO_VERTICAL_ANCHOR.CENTER
    for child_key, child_value in value.items():
        render_item_second(child_key, child_value)


def render_item_second(key, value):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    # 1、新增母版占位符
    placeholder1 = slide.placeholders[1]
    placeholder2 = slide.placeholders[2]
    picture = placeholder1.insert_picture('./pptx/static/bg.png')
    placeholder2.element.getparent().remove(placeholder2.element)
    
    # 2、设置占位符宽高
    picture.left = 0
    picture.top = 0
    picture.width = prs.slide_width
    picture.height = prs.slide_height
    # 3、添加文本框
    text_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(15), Cm(2))
    tf = text_box.text_frame
    tf.text = key
    # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    # tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    # 4、添加group形状组合
    group_shape = slide.shapes.add_group_shape()
    group_shape.height = prs.slide_height
    print(group_shape, '----------------', group_shape.shapes.build_freeform)
    group_shape.shapes.build_freeform(Inches(1), Inches(10), scale=3)
    group_shape.shapes.add_picture("./pptx/static/pc.png", Cm(1.5), Cm(5), Cm(10), Cm(8))

    if isinstance(value, list):
        # top = width = height = Inches(0.6)
        # oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(12), top, width, height)
        # fill = oval_shape.fill
        # fill.solid()
        # fill.fore_color.rgb = RGBColor(68, 84, 106)
        for child_idx, child_value in enumerate(value):
            text_box_common = slide.shapes.add_textbox(Cm(12), Cm(3) + child_idx * (prs.slide_height - Cm(3)*2)/(len(value) - 1), Cm(10), Cm(2))
            text_frame_common = text_box_common.text_frame
            text_frame_common.paragraphs[0].text = child_value if isinstance(child_value, str) else ""
            text_frame_common.paragraphs[0].font.bold = True
            text_frame_common.paragraphs[0].font.name = "宋体"
            text_frame_common.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
            text_frame_common.paragraphs[0].font.size = Pt(14)
            text_frame_common.word_wrap = True
            text_frame_common.paragraphs[0].word_wrap = True
            text_frame_common.paragraphs[0].vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame_common.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            print(type(child_value), '看看第三级的数据格式')
            if isinstance(child_value, list):
                for grand_idx, grand_value in enumerate(value):
                    p = text_frame_common.add_paragraph()
                    p.text = grand_value
                    p.font.bold = True
                    p.font.name = "宋体"
                    p.font.color.rgb = RGBColor(68, 84, 106)
                    p.font.size = Pt(10)
                    p.level = 2

    elif isinstance(value, dict):
        # 添加一个原型
        # left = top = width = height = Inches(0.6)
        # group_shape.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        # oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(12), top, width, height)
        # fill = oval_shape.fill
        # fill.solid()
        # fill.fore_color.rgb = RGBColor(68, 84, 106)
        child_idx = 0
        for child_key, child_value in value.items():
            child_idx = child_idx + 1
            print(value.keys(), 'ttttttttttttttttttttttt', len(value.keys()))
            text_box_common = slide.shapes.add_textbox(Cm(12), child_idx * (prs.slide_height - Cm(3) * 2) / (len(value.keys() - 1)),
                                                       Cm(10), Cm(2))
            text_frame_common = text_box_common.text_frame
            text_frame_common.word_wrap = True
            text_frame_common.paragraphs[0].text = child_value if isinstance(child_value, str) else ""
            text_frame_common.paragraphs[0].font.bold = True
            text_frame_common.paragraphs[0].font.name = "宋体"
            text_frame_common.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
            text_frame_common.paragraphs[0].font.size = Pt(14)
            text_frame_common.paragraphs[0].word_wrap = True
            text_frame_common.paragraphs[0].vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame_common.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            child_idx = child_idx + 1
            render_item_second(child_key, child_value)
    elif isinstance(value, str):
        text_box2 = slide.shapes.add_textbox(Cm(12), Cm(3), Cm(10), Cm(4))
        tf2 = text_box2.text_frame
        tf2.word_wrap = True
        p = tf2.add_paragraph()
        p.text = value
        p.font.bold = True
        p.font.name = "宋体"
        p.font.color.rgb = RGBColor(68, 84, 106)
        p.font.size = Pt(14)
        p.level = 2
        return
    else:
        return


def save_ppt(structured_list):
    for key, value in structured_list.items():
        render_item(key, value)
    prs.save('gpt测试ppt2.pptx')


data = get_gpt_data()
structured_data = clean_data(data)
save_ppt(structured_data)